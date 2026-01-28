# app/routes_admin.py

import logging
import gc
import re
import os
import tempfile
import threading
from io import BytesIO
import time
import uuid

import pandas as pd
from flask import (
    render_template, request, jsonify, session, redirect,
    url_for, flash, send_file
)
from functools import wraps
from datetime import datetime
from sqlalchemy import func
from werkzeug.utils import secure_filename
from dateutil import parser as date_parser

from .extensions import db
from .models import EmployeeRecord
from .aws_utils import (
    s3_manager,
    call_bedrock_with_context,
    synthesize_speech,
    transcribe_audio_from_s3,
    S3_VOICE_INPUT_PREFIX,
)
from .services import (
    read_employee_data_sample_from_rds,
    get_employee_stats_cached,
    get_chart_data_from_rds,
    get_advanced_insights_data,
    write_employee_data_to_s3,
    backup_current_data_in_s3,
    clean_employee_dataframe_chunked,
    create_knowledge_base_file,  # still available if needed elsewhere
    save_cleaned_data_to_rds,
    allowed_file,
    PANDAS_OPTIONS,
    invalidate_stats_cache,
    MAX_DISPLAY_RECORDS,
    REQUIRED_COLUMNS,
    CHUNK_SIZE,
    clear_employee_data_in_rds,
    save_cleaned_data_chunk_to_rds,
    write_employee_csv_string_to_s3,
    create_enhanced_knowledge_base_from_rds,
    get_existing_dedup_keys,
    filter_new_records_for_append,
    export_rds_to_csv_snapshot,
    _deduplicate_cleaned_dataframe,
    CSV_EXPORT_COLUMNS,
    generate_certificates_for_dataframe,
    generate_certificates_for_rds,
)

logger = logging.getLogger(__name__)

certificate_generation_lock = threading.Lock()
certificate_generation_state = {
    "running": False,
    "last_counts": None,
    "last_error": None,
    "last_finished_at": None,
}


def register_admin_routes(app):

    # --- Global 500 handler (JSON-safe for AJAX) --- #
    @app.errorhandler(500)
    def handle_internal_error(error):
        # If client expects JSON (e.g., fetch), return JSON instead of HTML
        if request.is_json or request.accept_mimetypes['application/json'] > request.accept_mimetypes['text/html']:
            logger.error(f"Global 500 error on path {request.path}: {error}")
            return jsonify({
                "success": False,
                "error": "Internal server error",
            }), 500
        # Otherwise, you can render a template or keep default
        logger.error(f"Global 500 error (HTML) on path {request.path}: {error}")
        return "Internal Server Error", 500

    # --- login decorator --- #
    def login_required(f):
        """Decorator to require login for admin functions (HTML + JSON safe)."""
        @wraps(f)
        def decorated_function(*args, **kwargs):
            if not session.get('logged_in'):
                # If it's an AJAX / JSON request, return JSON error instead of HTML
                if request.is_json or request.accept_mimetypes['application/json'] > request.accept_mimetypes['text/html']:
                    return jsonify({
                        'success': False,
                        'error': 'Login required',
                        'login_required': True
                    }), 401
                # Normal browser navigation → redirect to login page
                return redirect(url_for('admin_login'))
            return f(*args, **kwargs)
        return decorated_function

    SEARCHABLE_COLUMNS = {
        "assessment_id": ("Assessment ID", EmployeeRecord.assessment_id),
        "assessment_name": ("Assessment Name", EmployeeRecord.assessment_name),
        "name": ("Employee Name", EmployeeRecord.name),
        "email": ("Email", EmployeeRecord.email),
        "employee_id": ("Employee ID", EmployeeRecord.employee_id),
        "final_completion_date": ("Completion Date", EmployeeRecord.final_completion_date),
        "issuer": ("Issuer", EmployeeRecord.issuer),
        "level": ("Level", EmployeeRecord.level),
        "marks": ("Marks", EmployeeRecord.marks),
        "qualifier": ("Qualifier", EmployeeRecord.qualifier),
        "skill": ("Skill", EmployeeRecord.skill),
        "skill_id": ("Skill ID", EmployeeRecord.skill_id),
        "valid_till": ("Valid Till", EmployeeRecord.valid_till),
        "wipro_function": ("Wipro Function", EmployeeRecord.wipro_function),
    }
    EXACT_MATCH_COLUMNS = {"employee_id", "assessment_id", "skill_id"}

    def _parse_multi_values(raw_value: str) -> list:
        if not raw_value:
            return []
        tokens = re.split(r"[\s,;]+", raw_value.strip())
        return [token for token in tokens if token]

    def _normalize_export_filters(source: dict) -> dict:
        employee_id_raw = source.get("employee_id", "").strip()
        employee_ids = _parse_multi_values(employee_id_raw)
        filters = {
            "level": source.get("level", "").strip(),
            "issuer": source.get("issuer", "").strip(),
            "qualifier": source.get("qualifier", "").strip(),
            "skill": source.get("skill", "").strip(),
            "assessment_name": source.get("assessment_name", "").strip(),
            "employee_id": employee_ids[0] if len(employee_ids) == 1 else "",
            "employee_ids": employee_ids if len(employee_ids) > 1 else [],
            "assessment_id": source.get("assessment_id", "").strip(),
            "skill_id": source.get("skill_id", "").strip(),
            "name": source.get("name", "").strip(),
            "email": source.get("email", "").strip(),
            "marks": source.get("marks", "").strip(),
            "valid_till": source.get("valid_till", "").strip(),
            "wipro_function": source.get("wipro_function", "").strip(),
            "start_date": source.get("start_date", "").strip(),
            "end_date": source.get("end_date", "").strip(),
            "search_column": source.get("search_column", "").strip(),
            "search_value": source.get("search_value", "").strip(),
        }
        return {key: value for key, value in filters.items() if value}

    def _apply_export_filters(query, filters: dict):
        if filters.get("level"):
            query = query.filter(EmployeeRecord.level.ilike(f"%{filters['level']}%"))
        if filters.get("issuer"):
            query = query.filter(EmployeeRecord.issuer == filters["issuer"])
        if filters.get("qualifier"):
            query = query.filter(EmployeeRecord.qualifier.ilike(f"%{filters['qualifier']}%"))
        if filters.get("skill"):
            query = query.filter(EmployeeRecord.skill.ilike(f"%{filters['skill']}%"))
        if filters.get("assessment_name"):
            query = query.filter(EmployeeRecord.assessment_name.ilike(f"%{filters['assessment_name']}%"))
        if filters.get("employee_ids"):
            query = query.filter(EmployeeRecord.employee_id.in_(filters["employee_ids"]))
        if filters.get("employee_id"):
            employee_ids = _parse_multi_values(filters["employee_id"])
            if len(employee_ids) > 1:
                query = query.filter(EmployeeRecord.employee_id.in_(employee_ids))
            else:
                query = query.filter(EmployeeRecord.employee_id == filters["employee_id"])
        if filters.get("assessment_id"):
            query = query.filter(EmployeeRecord.assessment_id == filters["assessment_id"])
        if filters.get("skill_id"):
            query = query.filter(EmployeeRecord.skill_id == filters["skill_id"])
        if filters.get("name"):
            query = query.filter(EmployeeRecord.name.ilike(f"%{filters['name']}%"))
        if filters.get("email"):
            query = query.filter(EmployeeRecord.email.ilike(f"%{filters['email']}%"))
        if filters.get("marks"):
            query = query.filter(EmployeeRecord.marks.ilike(f"%{filters['marks']}%"))
        if filters.get("valid_till"):
            query = query.filter(EmployeeRecord.valid_till.ilike(f"%{filters['valid_till']}%"))
        if filters.get("wipro_function"):
            query = query.filter(EmployeeRecord.wipro_function.ilike(f"%{filters['wipro_function']}%"))
        if filters.get("search_column") and filters.get("search_value"):
            column_key = filters["search_column"]
            if column_key in SEARCHABLE_COLUMNS:
                column = SEARCHABLE_COLUMNS[column_key][1]
                if column_key in EXACT_MATCH_COLUMNS:
                    query = query.filter(column == filters["search_value"])
                else:
                    query = query.filter(column.ilike(f"%{filters['search_value']}%"))
        return query

    def _filter_records_by_date(records, filters: dict):
        start_date = filters.get("start_date")
        end_date = filters.get("end_date")
        if not start_date and not end_date:
            return records

        filtered = []
        start_dt = datetime.strptime(start_date, "%Y-%m-%d") if start_date else None
        end_dt = datetime.strptime(end_date, "%Y-%m-%d") if end_date else None

        for record in records:
            if not record.final_completion_date:
                continue
            try:
                record_date = date_parser.parse(record.final_completion_date)
            except Exception:
                continue

            if start_dt and record_date < start_dt:
                continue
            if end_dt and record_date > end_dt:
                continue
            filtered.append(record)
        return filtered

    def _build_explorer_payload(records, filters: dict) -> dict:
        total_records = len(records)
        unique_employees = len({
            record.employee_id
            for record in records
            if record.employee_id
        })
        completion_count = sum(1 for record in records if record.final_completion_date)
        completion_rate = round(
            (completion_count / total_records) * 100,
            1
        ) if total_records else 0

        level_counts = {}
        qualifier_counts = {}
        skill_counts = {}
        issuer_counts = {}
        function_counts = {}
        time_buckets = {}

        for record in records:
            level = (record.level or "Unspecified").strip() or "Unspecified"
            qualifier = (record.qualifier or "Unspecified").strip() or "Unspecified"
            skill = (record.skill or "Unspecified").strip() or "Unspecified"
            issuer = (record.issuer or "Unspecified").strip() or "Unspecified"
            function = (record.wipro_function or "Unspecified").strip() or "Unspecified"

            level_counts[level] = level_counts.get(level, 0) + 1
            qualifier_counts[qualifier] = qualifier_counts.get(qualifier, 0) + 1
            skill_counts[skill] = skill_counts.get(skill, 0) + 1
            issuer_counts[issuer] = issuer_counts.get(issuer, 0) + 1
            function_counts[function] = function_counts.get(function, 0) + 1

            if record.final_completion_date:
                try:
                    parsed_date = date_parser.parse(record.final_completion_date)
                except Exception:
                    parsed_date = None
                if parsed_date:
                    month_key = parsed_date.strftime("%Y-%m")
                    time_buckets[month_key] = time_buckets.get(month_key, 0) + 1

        def _top_item(counts: dict) -> dict:
            if not counts:
                return {"label": "N/A", "count": 0}
            label, count = max(counts.items(), key=lambda item: item[1])
            return {"label": label, "count": count}

        timeline_keys = sorted(time_buckets.keys())
        timeline_labels = [
            datetime.strptime(key, "%Y-%m").strftime("%b %Y")
            for key in timeline_keys
        ]
        timeline_counts = [time_buckets[key] for key in timeline_keys]

        records_preview = [
            {
                "employee_id": record.employee_id,
                "name": record.name,
                "email": record.email,
                "assessment_name": record.assessment_name,
                "level": record.level,
                "qualifier": record.qualifier,
                "final_completion_date": record.final_completion_date,
                "issuer": record.issuer,
                "skill": record.skill,
                "wipro_function": record.wipro_function,
            }
            for record in records[:MAX_DISPLAY_RECORDS]
        ]

        return {
            "totals": {
                "records": total_records,
                "unique_employees": unique_employees,
                "completion_rate": completion_rate,
                "certified": qualifier_counts.get("Certified", 0),
                "trained": qualifier_counts.get("Trained", 0),
            },
            "highlights": {
                "top_skill": _top_item(skill_counts),
                "top_issuer": _top_item(issuer_counts),
                "top_function": _top_item(function_counts),
            },
            "time_series": {
                "labels": timeline_labels,
                "counts": timeline_counts,
            },
            "level_distribution": {
                "labels": list(level_counts.keys()),
                "counts": list(level_counts.values()),
            },
            "qualifier_distribution": {
                "labels": list(qualifier_counts.keys()),
                "counts": list(qualifier_counts.values()),
            },
            "records": records_preview,
            "filters": filters,
        }

    def _build_filtered_summary(records: list) -> str:
        if not records:
            return "No matching records were found for the selected filters."

        total_records = len(records)
        unique_employees = len({
            record.employee_id
            for record in records
            if record.employee_id
        })
        level_counts = {}
        skill_counts = {}
        issuer_counts = {}
        for record in records:
            level = (record.level or "Unspecified").strip() or "Unspecified"
            skill = (record.skill or "Unspecified").strip() or "Unspecified"
            issuer = (record.issuer or "Unspecified").strip() or "Unspecified"
            level_counts[level] = level_counts.get(level, 0) + 1
            skill_counts[skill] = skill_counts.get(skill, 0) + 1
            issuer_counts[issuer] = issuer_counts.get(issuer, 0) + 1

        def _top_label(counts: dict) -> str:
            if not counts:
                return "N/A"
            return max(counts.items(), key=lambda item: item[1])[0]

        sample_rows = []
        for record in records[:5]:
            sample_rows.append(
                f"- {record.employee_id or 'Unknown ID'} | "
                f"{record.name or 'Unknown'} | "
                f"{record.assessment_name or 'Unknown Assessment'} | "
                f"{record.level or 'Unknown Level'} | "
                f"{record.qualifier or 'Unknown Qualifier'}"
            )

        return (
            "Filtered database summary:\n"
            f"- Total records: {total_records}\n"
            f"- Unique employees: {unique_employees}\n"
            f"- Top level: {_top_label(level_counts)}\n"
            f"- Top skill: {_top_label(skill_counts)}\n"
            f"- Top issuer: {_top_label(issuer_counts)}\n"
            "Sample records:\n"
            + "\n".join(sample_rows)
        )

    def _build_overall_summary() -> str:
        total_records = db.session.query(func.count(EmployeeRecord.id)).scalar() or 0
        unique_employees = db.session.query(
            func.count(func.distinct(EmployeeRecord.employee_id))
        ).filter(
            EmployeeRecord.employee_id.isnot(None),
            EmployeeRecord.employee_id != ''
        ).scalar() or 0
        date_range = db.session.query(
            func.min(EmployeeRecord.final_completion_date),
            func.max(EmployeeRecord.final_completion_date)
        ).one()
        earliest_date = date_range[0].strftime("%Y-%m-%d") if date_range[0] else "N/A"
        latest_date = date_range[1].strftime("%Y-%m-%d") if date_range[1] else "N/A"

        def _top_items(column, label, limit=5):
            rows = db.session.query(
                column,
                func.count(EmployeeRecord.id)
            ).filter(
                column.isnot(None),
                column != ''
            ).group_by(
                column
            ).order_by(
                func.count(EmployeeRecord.id).desc()
            ).limit(limit).all()
            if not rows:
                return f"- Top {label}: N/A"
            formatted = ", ".join([f"{value} ({count})" for value, count in rows])
            return f"- Top {label}: {formatted}"

        return (
            "Current database snapshot:\n"
            f"- Total records: {total_records}\n"
            f"- Unique employees: {unique_employees}\n"
            f"- Data range: {earliest_date} to {latest_date}\n"
            f"{_top_items(EmployeeRecord.issuer, 'issuers')}\n"
            f"{_top_items(EmployeeRecord.assessment_name, 'certifications')}\n"
            f"{_top_items(EmployeeRecord.skill, 'skills')}\n"
            f"{_top_items(EmployeeRecord.level, 'levels')}\n"
            f"{_top_items(EmployeeRecord.wipro_function, 'functions')}"
        )

    def _records_to_export_rows(records):
        rows = []
        for record in records:
            row = {
                label: getattr(record, column, "") or ""
                for column, label in CSV_EXPORT_COLUMNS
            }
            rows.append(row)
        return rows

    def _get_split_insights(issuer: str, split_by: str) -> dict:
        split_columns = {
            "wipro_function": ("Wipro Function", EmployeeRecord.wipro_function),
            "level": ("Level", EmployeeRecord.level),
            "qualifier": ("Qualifier", EmployeeRecord.qualifier),
            "skill": ("Skill", EmployeeRecord.skill),
        }
        if split_by not in split_columns:
            raise ValueError("Invalid split column.")

        split_label, split_column = split_columns[split_by]

        total_records = db.session.query(func.count(EmployeeRecord.id)).filter(
            EmployeeRecord.issuer == issuer
        ).scalar() or 0
        unique_employees = db.session.query(
            func.count(func.distinct(EmployeeRecord.employee_id))
        ).filter(
            EmployeeRecord.issuer == issuer,
            EmployeeRecord.employee_id.isnot(None),
            EmployeeRecord.employee_id != ''
        ).scalar() or 0
        function_count = db.session.query(
            func.count(func.distinct(EmployeeRecord.wipro_function))
        ).filter(
            EmployeeRecord.issuer == issuer,
            EmployeeRecord.wipro_function.isnot(None),
            EmployeeRecord.wipro_function != ''
        ).scalar() or 0
        skill_count = db.session.query(
            func.count(func.distinct(EmployeeRecord.skill))
        ).filter(
            EmployeeRecord.issuer == issuer,
            EmployeeRecord.skill.isnot(None),
            EmployeeRecord.skill != ''
        ).scalar() or 0

        certifications = db.session.query(
            EmployeeRecord.assessment_name,
            func.count(EmployeeRecord.id)
        ).filter(
            EmployeeRecord.issuer == issuer,
            EmployeeRecord.assessment_name.isnot(None),
            EmployeeRecord.assessment_name != ''
        ).group_by(
            EmployeeRecord.assessment_name
        ).order_by(
            func.count(EmployeeRecord.id).desc()
        ).limit(25).all()

        certification_data = [
            {"name": name, "count": count}
            for name, count in certifications
        ]

        split_rows = db.session.query(
            split_column,
            func.count(EmployeeRecord.id)
        ).filter(
            EmployeeRecord.issuer == issuer
        ).group_by(
            split_column
        ).order_by(
            func.count(EmployeeRecord.id).desc()
        ).limit(12).all()

        split_data = [
            {"label": value or "Unknown", "count": count}
            for value, count in split_rows
        ]

        skills = db.session.query(
            EmployeeRecord.skill,
            func.count(EmployeeRecord.id)
        ).filter(
            EmployeeRecord.issuer == issuer,
            EmployeeRecord.skill.isnot(None),
            EmployeeRecord.skill != ''
        ).group_by(
            EmployeeRecord.skill
        ).order_by(
            func.count(EmployeeRecord.id).desc()
        ).limit(15).all()

        skill_data = [
            {"name": name, "count": count}
            for name, count in skills
        ]

        top_certifications = [row[0] for row in certifications[:6]]
        top_functions = [
            row[0] or "Unknown"
            for row in db.session.query(
                EmployeeRecord.wipro_function,
                func.count(EmployeeRecord.id)
            ).filter(
                EmployeeRecord.issuer == issuer
            ).group_by(
                EmployeeRecord.wipro_function
            ).order_by(
                func.count(EmployeeRecord.id).desc()
            ).limit(6).all()
        ]

        matrix_rows = []
        if top_certifications and top_functions:
            matrix_rows = db.session.query(
                EmployeeRecord.assessment_name,
                EmployeeRecord.wipro_function,
                func.count(EmployeeRecord.id)
            ).filter(
                EmployeeRecord.issuer == issuer,
                EmployeeRecord.assessment_name.in_(top_certifications),
                EmployeeRecord.wipro_function.in_(top_functions)
            ).group_by(
                EmployeeRecord.assessment_name,
                EmployeeRecord.wipro_function
            ).all()

        matrix = {
            cert: {function: 0 for function in top_functions}
            for cert in top_certifications
        }
        for cert_name, function, count in matrix_rows:
            function_key = function or "Unknown"
            if cert_name in matrix and function_key in matrix[cert_name]:
                matrix[cert_name][function_key] = count

        top_cert = certification_data[0] if certification_data else {"name": "N/A", "count": 0}
        top_skill = skill_data[0] if skill_data else {"name": "N/A", "count": 0}
        top_split = split_data[0] if split_data else {"label": "N/A", "count": 0}

        return {
            "issuer": issuer,
            "totals": {
                "records": total_records,
                "employees": unique_employees,
                "certifications": len(certification_data),
                "functions": function_count,
                "skills": skill_count,
            },
            "certifications": certification_data,
            "skills": skill_data,
            "split_by": split_by,
            "split_label": split_label,
            "split_data": split_data,
            "matrix": {
                "certifications": top_certifications,
                "functions": top_functions,
                "data": matrix
            },
            "highlights": {
                "top_certification": top_cert,
                "top_skill": top_skill,
                "top_split": top_split,
            }
        }

    def _build_splits_pptx_report(split_payload: dict) -> BytesIO:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt

        prs = Presentation()
        issuer = split_payload["issuer"]
        split_label = split_payload["split_label"]
        totals = split_payload["totals"]
        highlights = split_payload["highlights"]

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Issuer Split Report"
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"Issuer: {issuer}\nSplit by: {split_label}"

        kpi_slide = prs.slides.add_slide(prs.slide_layouts[5])
        kpi_slide.shapes.title.text = "Issuer KPIs"
        kpi_box = kpi_slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.5), Inches(3))
        kpi_frame = kpi_box.text_frame
        kpi_frame.word_wrap = True
        kpi_frame.text = f"Total Records: {totals['records']}"
        for line in [
            f"Unique Employees: {totals['employees']}",
            f"Certifications: {totals['certifications']}",
            f"Functions Covered: {totals['functions']}",
            f"Skills Covered: {totals['skills']}",
            f"Top Certification: {highlights['top_certification']['name']} ({highlights['top_certification']['count']})",
            f"Top Skill: {highlights['top_skill']['name']} ({highlights['top_skill']['count']})",
            f"Top {split_label}: {highlights['top_split']['label']} ({highlights['top_split']['count']})",
        ]:
            paragraph = kpi_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(18)

        split_slide = prs.slides.add_slide(prs.slide_layouts[5])
        split_slide.shapes.title.text = f"{split_label} Split"
        chart_data = ChartData()
        split_data = split_payload["split_data"]
        chart_data.categories = [item["label"] for item in split_data] or ["No Data"]
        chart_data.add_series(
            "Completions",
            [item["count"] for item in split_data] or [0]
        )
        split_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8),
            Inches(1.6),
            Inches(8.6),
            Inches(4.5),
            chart_data,
        )

        insight_slide = prs.slides.add_slide(prs.slide_layouts[5])
        insight_slide.shapes.title.text = "Top Certifications & Skills"
        text_box = insight_slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.5), Inches(4))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = "Certifications:"
        for item in split_payload["certifications"][:8]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = f"- {item['name']}: {item['count']}"
            paragraph.level = 1
        text_frame.add_paragraph().text = ""
        skills_header = text_frame.add_paragraph()
        skills_header.text = "Skills:"
        for item in split_payload["skills"][:8]:
            paragraph = text_frame.add_paragraph()
            paragraph.text = f"- {item['name']}: {item['count']}"
            paragraph.level = 1

        matrix = split_payload["matrix"]
        if matrix["certifications"] and matrix["functions"]:
            matrix_slide = prs.slides.add_slide(prs.slide_layouts[5])
            matrix_slide.shapes.title.text = "Top Certifications by Function"
            rows = len(matrix["certifications"]) + 1
            cols = len(matrix["functions"]) + 1
            table = matrix_slide.shapes.add_table(
                rows, cols, Inches(0.4), Inches(1.4), Inches(9.1), Inches(4.8)
            ).table
            table.cell(0, 0).text = "Certification"
            for idx, fn in enumerate(matrix["functions"], start=1):
                table.cell(0, idx).text = fn
            for row_idx, cert in enumerate(matrix["certifications"], start=1):
                table.cell(row_idx, 0).text = cert
                for col_idx, fn in enumerate(matrix["functions"], start=1):
                    table.cell(row_idx, col_idx).text = str(
                        matrix["data"].get(cert, {}).get(fn, 0)
                    )

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def _build_pptx_report(records, filters: dict):
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = "AI Academy Advanced Insights"
        subtitle = title_slide.placeholders[1]
        filter_lines = [f"{key.replace('_', ' ').title()}: {value}" for key, value in filters.items()]
        subtitle.text = "Filters applied:\n" + ("\n".join(filter_lines) if filter_lines else "All records")

        # KPI slide
        kpi_slide = prs.slides.add_slide(prs.slide_layouts[5])
        kpi_slide.shapes.title.text = "KPI Summary"
        total_records = len(records)
        unique_employees = len({rec.employee_id for rec in records if rec.employee_id})
        level_2_plus = sum(
            1 for rec in records
            if rec.level and ("Level 2" in rec.level or "Level 3" in rec.level)
        )
        kpi_box = kpi_slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(3))
        text_frame = kpi_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = f"Total Records: {total_records}"
        for line in [
            f"Unique Employees: {unique_employees}",
            f"Level 2+ Achievers: {level_2_plus}",
        ]:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)

        # Level distribution chart
        level_counts = {}
        for rec in records:
            level = (rec.level or "Unknown").strip()
            level_counts[level] = level_counts.get(level, 0) + 1
        chart_data = ChartData()
        chart_data.categories = list(level_counts.keys()) or ["No Data"]
        chart_data.add_series("Level Distribution", list(level_counts.values()) or [0])
        chart_slide = prs.slides.add_slide(prs.slide_layouts[5])
        chart_slide.shapes.title.text = "Level Distribution"
        chart_slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8),
            Inches(1.6),
            Inches(8.6),
            Inches(4.5),
            chart_data,
        )

        # Issuer distribution chart
        issuer_counts = {}
        for rec in records:
            issuer = (rec.issuer or "Unknown").strip()
            issuer_counts[issuer] = issuer_counts.get(issuer, 0) + 1
        issuer_chart_data = ChartData()
        top_issuers = sorted(issuer_counts.items(), key=lambda x: x[1], reverse=True)[:10]
        issuer_chart_data.categories = [item[0] for item in top_issuers] or ["No Data"]
        issuer_chart_data.add_series("Issuer Distribution", [item[1] for item in top_issuers] or [0])
        issuer_slide = prs.slides.add_slide(prs.slide_layouts[5])
        issuer_slide.shapes.title.text = "Top Issuers"
        issuer_slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.8),
            Inches(1.6),
            Inches(8.6),
            Inches(4.5),
            issuer_chart_data,
        )

        employee_ids = list(filters.get("employee_ids", []))
        if filters.get("employee_id"):
            employee_ids.append(filters["employee_id"])

        def _employee_snapshot(emp_id: str, emp_records: list) -> dict:
            level_counts = {}
            skill_counts = {}
            issuer_counts = {}
            completion_dates = []
            for rec in emp_records:
                if rec.level:
                    level_counts[rec.level] = level_counts.get(rec.level, 0) + 1
                if rec.skill:
                    skill_counts[rec.skill] = skill_counts.get(rec.skill, 0) + 1
                if rec.issuer:
                    issuer_counts[rec.issuer] = issuer_counts.get(rec.issuer, 0) + 1
                if rec.final_completion_date:
                    try:
                        completion_dates.append(date_parser.parse(rec.final_completion_date))
                    except Exception:
                        continue

            def _top_item(counts: dict) -> str:
                if not counts:
                    return "N/A"
                return max(counts.items(), key=lambda item: item[1])[0]

            latest_completion = max(completion_dates).strftime("%Y-%m-%d") if completion_dates else "N/A"
            return {
                "employee_id": emp_id,
                "records": len(emp_records),
                "top_level": _top_item(level_counts),
                "top_skill": _top_item(skill_counts),
                "top_issuer": _top_item(issuer_counts),
                "latest_completion": latest_completion,
            }

        employee_summaries = []
        for emp_id in employee_ids:
            emp_records = [rec for rec in records if rec.employee_id == emp_id]
            if not emp_records:
                continue
            summary = _employee_snapshot(emp_id, emp_records)
            employee_summaries.append(summary)
            emp_slide = prs.slides.add_slide(prs.slide_layouts[5])
            emp_slide.shapes.title.text = f"Employee {emp_id} Snapshot"
            emp_box = emp_slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.5), Inches(3.5))
            emp_frame = emp_box.text_frame
            emp_frame.word_wrap = True
            emp_frame.text = f"Total Records: {summary['records']}"
            for line in [
                f"Top Level: {summary['top_level']}",
                f"Top Skill: {summary['top_skill']}",
                f"Top Issuer: {summary['top_issuer']}",
                f"Latest Completion: {summary['latest_completion']}",
            ]:
                paragraph = emp_frame.add_paragraph()
                paragraph.text = line
                paragraph.font.size = Pt(18)

        if len(employee_summaries) > 1:
            compare_slide = prs.slides.add_slide(prs.slide_layouts[5])
            compare_slide.shapes.title.text = "Employee Comparison"
            rows = len(employee_summaries) + 1
            cols = 5
            table = compare_slide.shapes.add_table(
                rows, cols, Inches(0.4), Inches(1.4), Inches(9.1), Inches(4.8)
            ).table
            table.cell(0, 0).text = "Employee ID"
            table.cell(0, 1).text = "Records"
            table.cell(0, 2).text = "Top Level"
            table.cell(0, 3).text = "Top Skill"
            table.cell(0, 4).text = "Latest Completion"
            for row_idx, summary in enumerate(employee_summaries, start=1):
                table.cell(row_idx, 0).text = summary["employee_id"]
                table.cell(row_idx, 1).text = str(summary["records"])
                table.cell(row_idx, 2).text = summary["top_level"]
                table.cell(row_idx, 3).text = summary["top_skill"]
                table.cell(row_idx, 4).text = summary["latest_completion"]

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def _extract_export_filters_from_message(message: str) -> dict:
        filters = {}
        normalized_message = message.lower()
        level_match = re.search(r"level\s*(\d+)", message, re.IGNORECASE)
        if level_match:
            filters["level"] = f"Level {level_match.group(1)}"

        if "certified" in normalized_message:
            filters["qualifier"] = "Certified"
        elif "trained" in normalized_message:
            filters["qualifier"] = "Trained"

        employee_block = re.search(
            r"(employee\s*ids?|emp\s*ids?)\s*[:#]?\s*([0-9,\s]+)",
            normalized_message
        )
        if employee_block:
            employee_ids = re.findall(r"\d{4,}", employee_block.group(2))
            if len(employee_ids) > 1:
                filters["employee_ids"] = employee_ids
            elif employee_ids:
                filters["employee_id"] = employee_ids[0]
        else:
            employee_id_match = re.search(
                r"(employee\s*id|emp\s*id)\s*[:#]?\s*(\d{4,})",
                normalized_message
            )
            if employee_id_match:
                filters["employee_id"] = employee_id_match.group(2)

        email_match = re.search(r"[\w\.-]+@[\w\.-]+\.\w+", message)
        if email_match:
            filters["email"] = email_match.group(0)

        assessment_id_match = re.search(
            r"assessment\s*id\s*[:#]?\s*([\w-]+)",
            normalized_message
        )
        if assessment_id_match:
            filters["assessment_id"] = assessment_id_match.group(1)

        skill_id_match = re.search(
            r"skill\s*id\s*[:#]?\s*([\w-]+)",
            normalized_message
        )
        if skill_id_match:
            filters["skill_id"] = skill_id_match.group(1)

        marks_match = re.search(
            r"marks\s*[:#]?\s*([\w-]+)",
            normalized_message
        )
        if marks_match:
            filters["marks"] = marks_match.group(1)

        valid_till_match = re.search(
            r"valid\s*till\s*[:#]?\s*([a-zA-Z0-9\-/]+)",
            normalized_message
        )
        if valid_till_match:
            filters["valid_till"] = valid_till_match.group(1)

        name_match = re.search(
            r"(employee\s*name|name)\s*(?:is|:)?\s*([a-zA-Z.'\s-]{2,})",
            message,
            re.IGNORECASE
        )
        if name_match:
            filters["name"] = name_match.group(2).strip()

        issuer_candidates = db.session.query(func.distinct(EmployeeRecord.issuer)).filter(
            EmployeeRecord.issuer.isnot(None),
            EmployeeRecord.issuer != ''
        ).all()
        issuers = sorted(
            [issuer[0] for issuer in issuer_candidates if issuer[0]],
            key=len,
            reverse=True
        )
        for issuer in issuers:
            if issuer.lower() in message.lower():
                filters["issuer"] = issuer
                break

        skill_candidates = db.session.query(func.distinct(EmployeeRecord.skill)).filter(
            EmployeeRecord.skill.isnot(None),
            EmployeeRecord.skill != ''
        ).all()
        skills = sorted(
            [skill[0] for skill in skill_candidates if skill[0]],
            key=len,
            reverse=True
        )
        for skill in skills:
            if skill.lower() in message.lower():
                filters["skill"] = skill
                break

        assessment_candidates = db.session.query(func.distinct(EmployeeRecord.assessment_name)).filter(
            EmployeeRecord.assessment_name.isnot(None),
            EmployeeRecord.assessment_name != ''
        ).all()
        assessments = sorted(
            [assessment[0] for assessment in assessment_candidates if assessment[0]],
            key=len,
            reverse=True
        )
        for assessment in assessments:
            if assessment.lower() in message.lower():
                filters["assessment_name"] = assessment
                break

        function_candidates = db.session.query(func.distinct(EmployeeRecord.wipro_function)).filter(
            EmployeeRecord.wipro_function.isnot(None),
            EmployeeRecord.wipro_function != ''
        ).all()
        functions = sorted(
            [function[0] for function in function_candidates if function[0]],
            key=len,
            reverse=True
        )
        for function in functions:
            if function.lower() in message.lower():
                filters["wipro_function"] = function
                break

        field_match = re.search(
            r"(?:field|column)\s+(?P<column>[a-zA-Z_\s]+?)\s*(?:is|=|:)\s*(?P<value>[^,]+)",
            message,
            re.IGNORECASE
        )
        if field_match:
            column_raw = field_match.group("column").strip().lower().replace("_", " ")
            value = field_match.group("value").strip()
            column_lookup = {
                "employee id": "employee_id",
                "emp id": "employee_id",
                "assessment id": "assessment_id",
                "assessment name": "assessment_name",
                "employee name": "name",
                "name": "name",
                "email": "email",
                "issuer": "issuer",
                "level": "level",
                "qualifier": "qualifier",
                "skill": "skill",
                "skill id": "skill_id",
                "valid till": "valid_till",
                "wipro function": "wipro_function",
                "marks": "marks",
            }
            column_key = column_lookup.get(column_raw)
            if column_key:
                if column_key == "level" and value.isdigit():
                    filters[column_key] = f"Level {value}"
                elif column_key == "employee_id":
                    employee_ids = _parse_multi_values(value)
                    if len(employee_ids) > 1:
                        filters["employee_ids"] = employee_ids
                    elif employee_ids:
                        filters["employee_id"] = employee_ids[0]
                else:
                    filters[column_key] = value

        explicit_field_patterns = {
            "assessment_id": r"assessment\s*id\s*[:=]\s*([\\w-]+)",
            "assessment_name": r"assessment\s*name\s*[:=]\s*([\\w\\s\\-_.]+)",
            "name": r"employee\s*name\s*[:=]\s*([a-zA-Z.'\\s-]{2,})",
            "email": r"email\s*[:=]\s*([\\w\\.-]+@[\\w\\.-]+\\.[a-zA-Z]+)",
            "skill": r"skill\s*[:=]\s*([\\w\\s\\-_.]+)",
            "skill_id": r"skill\s*id\s*[:=]\s*([\\w-]+)",
            "issuer": r"issuer\s*[:=]\s*([\\w\\s\\-_.]+)",
            "level": r"level\s*[:=]\s*(\\d+)",
            "qualifier": r"qualifier\s*[:=]\s*([\\w\\s\\-_.]+)",
            "wipro_function": r"wipro\\s*function\s*[:=]\s*([\\w\\s\\-_.]+)",
            "employee_id": r"employee\\s*id\\s*[:=]\s*(\\d+)",
            "marks": r"marks\\s*[:=]\\s*([\\w\\s\\-_.]+)",
            "valid_till": r"valid\\s*till\\s*[:=]\\s*([\\w\\s\\-_.]+)",
        }

        for key, pattern in explicit_field_patterns.items():
            if key in filters:
                continue
            match = re.search(pattern, message, re.IGNORECASE)
            if match:
                value = match.group(1).strip()
                if key == "level" and value.isdigit():
                    filters[key] = f"Level {value}"
                else:
                    filters[key] = value

        date_filters = _extract_date_filters_from_message(message)
        filters.update(date_filters)

        return filters

    def _extract_date_filters_from_message(message: str) -> dict:
        patterns = [
            r"between\s+(?P<start>[^,]+?)\s+and\s+(?P<end>[^,]+)",
            r"from\s+(?P<start>[^,]+?)\s+to\s+(?P<end>[^,]+)",
        ]
        for pattern in patterns:
            match = re.search(pattern, message, re.IGNORECASE)
            if match:
                start_date = _safe_parse_date(match.group("start"))
                end_date = _safe_parse_date(match.group("end"))
                return _build_date_filter_payload(start_date, end_date)

        explicit_dates = re.findall(r"\b\d{4}-\d{2}-\d{2}\b", message)
        if len(explicit_dates) >= 2:
            return _build_date_filter_payload(explicit_dates[0], explicit_dates[1])

        start_match = re.search(
            r"(after|since|from)\s+(?P<start>[^,]+)",
            message,
            re.IGNORECASE
        )
        end_match = re.search(
            r"(before|until|up to)\s+(?P<end>[^,]+)",
            message,
            re.IGNORECASE
        )
        start_date = _safe_parse_date(start_match.group("start")) if start_match else None
        end_date = _safe_parse_date(end_match.group("end")) if end_match else None
        return _build_date_filter_payload(start_date, end_date)

    def _safe_parse_date(raw_value: str):
        if not raw_value:
            return None
        try:
            parsed = date_parser.parse(raw_value, fuzzy=True)
            return parsed.date().isoformat()
        except Exception:
            return None

    def _build_date_filter_payload(start_date, end_date) -> dict:
        payload = {}
        if start_date:
            payload["start_date"] = start_date
        if end_date:
            payload["end_date"] = end_date
        return payload

    # --- Auth routes --- #

    @app.route('/admin/login', methods=['GET', 'POST'])
    def admin_login():
        """Admin login."""
        try:
            if request.method == 'POST':
                username = request.form.get('username', '').strip()
                password = request.form.get('password', '').strip()
                logger.info(f"Login attempt - Username: {username}")

                if not username or not password:
                    flash('Please enter both username and password', 'error')
                    return render_template('login.html')

                admin_username = app.config.get('ADMIN_USERNAME', 'admin')
                admin_password = app.config.get('ADMIN_PASSWORD', 'magicalacademy123')

                if username == admin_username and password == admin_password:
                    session['logged_in'] = True
                    session['username'] = username
                    logger.info(f"✅ Successful login: {username}")
                    flash('Successfully logged in!', 'success')
                    return redirect(url_for('admin_panel'))
                else:
                    logger.warning(f"❌ Failed login for user: {username}")
                    flash('Invalid username or password', 'error')
                    return render_template('login.html')

            if session.get('logged_in'):
                return redirect(url_for('admin_panel'))

            return render_template('login.html')
        except Exception as e:
            logger.error(f"❌ Login error: {e}")
            flash('Login system error', 'error')
            return render_template('login.html')

    @app.route('/admin/logout')
    def admin_logout():
        """Admin logout."""
        try:
            session.clear()
            flash('Successfully logged out', 'success')
            return redirect(url_for('admin_login'))
        except Exception as e:
            logger.error(f"❌ Logout error: {e}")
            return redirect(url_for('admin_login'))

    # --- Admin panel and pages --- #

    @app.route('/admin')
    @login_required
    def admin_panel():
        """RDS-optimized admin panel."""
        try:
            logger.info(f"📊 Admin panel accessed by: {session.get('username')}")
            employees = read_employee_data_sample_from_rds()
            if len(employees) == MAX_DISPLAY_RECORDS:
                flash(
                    f'Showing first {MAX_DISPLAY_RECORDS} records for performance. '
                    f'Use search or database browser to find specific employees.',
                    'info'
                )
            stats = get_employee_stats_cached()
            return render_template(
                'admin.html',
                employees=employees,
                username=session.get('username'),
                stats=stats
            )
        except Exception as e:
            logger.error(f"❌ Admin panel error: {e}")
            return f"Admin panel error: {str(e)}", 500

    @app.route('/admin/upload', methods=['GET'])
    @login_required
    def upload_csv():
        """Upload page."""
        try:
            s3_enabled = s3_manager.s3_status.startswith('Ready')
            return render_template(
                'upload.html',
                username=session.get('username'),
                s3_enabled=s3_enabled
            )
        except Exception as e:
            logger.error(f"❌ Upload page error: {e}")
            return f"Upload page error: {str(e)}", 500

    @app.route('/admin/chat')
    @login_required
    def admin_chat():
        """AI Chat page."""
        try:
            return render_template('chat.html', username=session.get('username'))
        except Exception as e:
            logger.error(f"❌ Chat page error: {e}")
            return f"Chat page error: {str(e)}", 500

    @app.route('/admin/agentic')
    @login_required
    def admin_agentic():
        """Agentic AI reporting page."""
        try:
            return render_template('agentic.html', username=session.get('username'))
        except Exception as e:
            logger.error(f"❌ Agentic page error: {e}")
            return f"Agentic page error: {str(e)}", 500

    @app.route('/admin/insights')
    @login_required
    def admin_insights():
        """Advanced insights page for executive dashboards."""
        try:
            stats = get_employee_stats_cached()
            insights = get_advanced_insights_data()
            return render_template(
                'insights.html',
                username=session.get('username'),
                stats=stats,
                insights=insights
            )
        except Exception as e:
            logger.error(f"❌ Insights page error: {e}")
            return f"Insights page error: {str(e)}", 500

    @app.route('/admin/explorer')
    @login_required
    def admin_explorer():
        """Professional search and analytics explorer."""
        try:
            columns = [
                {"value": key, "label": label}
                for key, (label, _) in SEARCHABLE_COLUMNS.items()
            ]
            return render_template(
                'explorer.html',
                username=session.get('username'),
                columns=columns
            )
        except Exception as e:
            logger.error(f"❌ Explorer page error: {e}")
            return f"Explorer page error: {str(e)}", 500

    @app.route('/admin/explorer-data')
    @login_required
    def admin_explorer_data():
        """JSON data for the explorer search page."""
        try:
            filters = _normalize_export_filters(request.args)
            query = _apply_export_filters(db.session.query(EmployeeRecord), filters)
            records = query.all()
            records = _filter_records_by_date(records, filters)
            payload = _build_explorer_payload(records, filters)
            return jsonify({
                "success": True,
                **payload
            })
        except Exception as e:
            logger.error(f"Explorer data error: {e}")
            return jsonify({"success": False, "error": str(e)}), 500

    @app.route('/admin/splits')
    @login_required
    def admin_splits():
        """Issuer and certification split explorer."""
        try:
            issuers = [
                row[0]
                for row in db.session.query(EmployeeRecord.issuer)
                .filter(EmployeeRecord.issuer.isnot(None), EmployeeRecord.issuer != '')
                .distinct()
                .order_by(EmployeeRecord.issuer)
                .all()
            ]
            return render_template(
                'splits.html',
                username=session.get('username'),
                issuers=issuers
            )
        except Exception as e:
            logger.error(f"❌ Splits page error: {e}")
            return f"Splits page error: {str(e)}", 500

    @app.route('/admin/splits-data')
    @login_required
    def admin_splits_data():
        """JSON data for issuer splits (certifications + function/level split)."""
        issuer = request.args.get('issuer', '').strip()
        split_by = request.args.get('split_by', 'wipro_function').strip()
        if not issuer:
            return jsonify({"success": False, "error": "Issuer is required."}), 400

        try:
            payload = _get_split_insights(issuer, split_by)
        except ValueError as exc:
            return jsonify({"success": False, "error": str(exc)}), 400

        return jsonify({
            "success": True,
            **payload
        })

    @app.route('/admin/splits-export')
    @login_required
    def admin_splits_export():
        """Export issuer split data to Excel or PPTX."""
        issuer = request.args.get('issuer', '').strip()
        split_by = request.args.get('split_by', 'wipro_function').strip()
        export_format = request.args.get('format', 'xlsx').lower()
        if not issuer:
            return jsonify({"success": False, "error": "Issuer is required."}), 400

        try:
            payload = _get_split_insights(issuer, split_by)
        except ValueError as exc:
            return jsonify({"success": False, "error": str(exc)}), 400

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        if export_format == 'pptx':
            pptx_stream = _build_splits_pptx_report(payload)
            filename = f"issuer_split_{timestamp}.pptx"
            return send_file(
                pptx_stream,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=filename
            )

        output = BytesIO()
        summary = {
            "Issuer": payload["issuer"],
            "Split By": payload["split_label"],
            "Total Records": payload["totals"]["records"],
            "Unique Employees": payload["totals"]["employees"],
            "Certifications": payload["totals"]["certifications"],
            "Functions Covered": payload["totals"]["functions"],
            "Skills Covered": payload["totals"]["skills"],
            "Top Certification": payload["highlights"]["top_certification"]["name"],
            "Top Certification Count": payload["highlights"]["top_certification"]["count"],
            "Top Skill": payload["highlights"]["top_skill"]["name"],
            "Top Skill Count": payload["highlights"]["top_skill"]["count"],
            "Top Split Dimension": payload["highlights"]["top_split"]["label"],
            "Top Split Count": payload["highlights"]["top_split"]["count"],
        }

        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            pd.DataFrame([summary]).to_excel(writer, index=False, sheet_name='Summary')
            pd.DataFrame(payload["certifications"]).to_excel(
                writer, index=False, sheet_name='Certifications'
            )
            pd.DataFrame(payload["split_data"]).to_excel(
                writer, index=False, sheet_name='Split'
            )
            pd.DataFrame(payload["skills"]).to_excel(
                writer, index=False, sheet_name='Skills'
            )
            matrix = payload["matrix"]
            if matrix["certifications"] and matrix["functions"]:
                matrix_df = pd.DataFrame(matrix["data"]).T
                matrix_df = matrix_df.reindex(matrix["certifications"])
                matrix_df = matrix_df[matrix["functions"]]
                matrix_df.to_excel(writer, sheet_name='Matrix')

        output.seek(0)
        filename = f"issuer_split_{timestamp}.xlsx"
        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=filename
        )

    # --- Chat API --- #

    @app.route('/admin/chat-api', methods=['POST'])
    @login_required
    def chat_api():
        """Smart chat API with knowledge base + database queries (via Bedrock)."""
        try:
            data = request.get_json()
            user_message = data.get('message', '').strip()
            if not user_message:
                return jsonify({"success": False, "error": "Please enter a message"})

            kb_content, error = s3_manager.read_file_from_s3(
                'knowledge-base/employee_data.txt'
            )
            kb_text = kb_content.decode('utf-8') if kb_content else ""

            employee_ids = re.findall(r'\b\d{5,}\b', user_message)
            employee_data_text = ""
            if employee_ids:
                for emp_id in employee_ids[:5]:
                    records = EmployeeRecord.query.filter_by(
                        employee_id=emp_id
                    ).all()
                    if records:
                        employee_data_text += f"\n\nEmployee ID {emp_id} Records:\n"
                        for rec in records:
                            employee_data_text += (
                                f"- Level: {rec.level}, Qualifier: {rec.qualifier}, "
                                f"Issuer: {rec.issuer}, Valid Till: {rec.valid_till}\n"
                            )

            export_filters = _extract_export_filters_from_message(user_message)
            filtered_summary = ""
            if export_filters:
                query = _apply_export_filters(db.session.query(EmployeeRecord), export_filters)
                filtered_records = query.all()
                filtered_records = _filter_records_by_date(filtered_records, export_filters)
                filtered_summary = _build_filtered_summary(filtered_records)
            overall_summary = _build_overall_summary()

            employee_block = (
                f"SPECIFIC EMPLOYEE DATA:{employee_data_text}"
                if employee_data_text
                else ""
            )
            filtered_block = (
                f"FILTERED DATABASE SNAPSHOT:\n{filtered_summary}"
                if filtered_summary
                else ""
            )
            overall_block = (
                f"DATABASE SNAPSHOT:\n{overall_summary}"
                if overall_summary
                else ""
            )
            kb_block = (
                f"STATISTICAL DATA:\n{kb_text}"
                if kb_text
                else "STATISTICAL DATA:\nKnowledge base not available."
            )
            system_prompt = f"""You are an AI assistant for the AI Academy Tracker system.

{kb_block}

{employee_block}
{filtered_block}
{overall_block}

Instructions:
- Use the statistical data above for general queries about levels, issuers, skills, certifications, qualifiers, and functions.
- Use the specific employee data (if provided) for individual employee queries.
- Use the database snapshot to ground high-level questions when no filters are provided.
- The database includes: assessment_id, assessment_name, name, email, employee_id, final_completion_date, issuer, level, marks, qualifier, skill, skill_id, valid_till, wipro_function.
- If filter context is provided, summarize it with KPIs and grounded insights.
- If asked about a specific employee not in the data, say you need to query the database.
- Provide production-grade analysis with bullet points, KPIs, and next-step insights.
- Offer segmentation options by issuer, certification (assessment), skill, level, qualifier, and Wipro function.
- Highlight standout performers, gaps, and potential actions (enablement, training, hiring, or certification drives).
- When a user mentions exporting, summarize the filters you interpreted from their request.
- If the question is ambiguous, ask 1-2 clarifying questions and suggest the next best filters.

Answer the user's question below."""

            response_text = call_bedrock_with_context(user_message, system_prompt)

            export_requested = bool(
                re.search(r"\b(export|download|excel|xlsx|pptx|powerpoint|report)\b", user_message, re.IGNORECASE)
            )
            export_ready = export_requested or bool(export_filters)
            return jsonify({
                "success": True,
                "response": response_text,
                # Updated to reflect the AWS-native model
                "model_used": "Amazon Nova Pro (Bedrock)",
                "export_ready": export_ready,
                "export_filters": export_filters
            })
        except Exception as e:
            logger.error(f"Chat API error: {e}")
            return jsonify({"success": False, "error": str(e)})

    # --- Voice chat endpoints --- #

    @app.route('/admin/voice/transcribe', methods=['POST'])
    @login_required
    def voice_transcribe():
        """Transcribe a voice prompt using Amazon Transcribe."""
        try:
            audio_file = request.files.get('audio')
            if not audio_file:
                return jsonify({"success": False, "error": "No audio file uploaded"}), 400

            filename = secure_filename(audio_file.filename or "voice.webm")
            extension = filename.rsplit('.', 1)[-1].lower()
            if extension not in {'webm', 'ogg', 'wav', 'mp3', 'mp4', 'flac'}:
                return jsonify({"success": False, "error": "Unsupported audio format"}), 400

            s3_key = f"{S3_VOICE_INPUT_PREFIX}{uuid.uuid4()}.{extension}"
            audio_bytes = audio_file.read()
            success, error = s3_manager.write_file_to_s3(
                s3_key,
                audio_bytes,
                content_type=audio_file.mimetype or "audio/webm"
            )
            if not success:
                return jsonify({"success": False, "error": error or "Failed to store audio"}), 500

            transcript, error = transcribe_audio_from_s3(
                s3_key,
                media_format=extension
            )
            s3_manager.delete_file(s3_key)
            if error:
                return jsonify({"success": False, "error": error}), 500

            return jsonify({"success": True, "transcript": transcript})
        except Exception as e:
            logger.error(f"Voice transcription error: {e}")
            return jsonify({"success": False, "error": str(e)}), 500

    @app.route('/admin/voice/synthesize', methods=['POST'])
    @login_required
    def voice_synthesize():
        """Synthesize a response using Amazon Polly."""
        try:
            data = request.get_json()
            text = data.get("text", "").strip() if data else ""
            if not text:
                return jsonify({"success": False, "error": "Text is required"}), 400

            audio_bytes, error = synthesize_speech(text)
            if error or not audio_bytes:
                return jsonify({"success": False, "error": error or "Polly error"}), 500

            return send_file(
                BytesIO(audio_bytes),
                mimetype="audio/mpeg",
                as_attachment=False,
                download_name="response.mp3"
            )
        except Exception as e:
            logger.error(f"Voice synthesis error: {e}")
            return jsonify({"success": False, "error": str(e)}), 500

    # --- Change password (placeholder) --- #

    def load_admin_password() -> str:
        """Placeholder – adapt if you had a custom implementation."""
        return app.config.get('ADMIN_PASSWORD', 'magicalacademy123')

    def save_new_password(new_password: str) -> bool:
        """Placeholder – adapt to persist password as you did before."""
        app.config['ADMIN_PASSWORD'] = new_password
        return True

    @app.route('/admin/change-password', methods=['GET', 'POST'])
    @login_required
    def change_password():
        """Change admin password."""
        try:
            if request.method == 'POST':
                current_password = request.form.get('current_password', '').strip()
                new_password = request.form.get('new_password', '').strip()
                confirm_password = request.form.get('confirm_password', '').strip()

                if not all([current_password, new_password, confirm_password]):
                    flash('All password fields are required', 'error')
                    return render_template(
                        'change_password.html',
                        username=session.get('username')
                    )

                stored_password = load_admin_password()

                if current_password != stored_password:
                    flash('Current password is incorrect', 'error')
                elif len(new_password) < 6:
                    flash(
                        'New password must be at least 6 characters long',
                        'error'
                    )
                elif new_password != confirm_password:
                    flash(
                        'New password and confirmation do not match',
                        'error'
                    )
                else:
                    if save_new_password(new_password):
                        flash('Password updated successfully!', 'success')
                        return redirect(url_for('admin_panel'))
                    else:
                        flash('Error updating password. Please try again.', 'error')

            return render_template(
                'change_password.html',
                username=session.get('username')
            )
        except Exception as e:
            logger.error(f"❌ Change password error: {e}")
            return f"Password change service error: {str(e)}", 500

    # --- Upload URL & processing --- #

    @app.route('/admin/upload-url', methods=['POST'])
    @login_required
    def get_upload_url():
        """Generate presigned URL for S3 upload."""
        try:
            data = request.get_json()
            if not data:
                return jsonify({'error': 'Invalid request format'})

            filename = data.get('filename', '').strip()
            content_type = data.get('content_type', '').strip()

            if not filename or not content_type:
                return jsonify({'error': 'Filename and content type are required'})

            if not allowed_file(filename):
                return jsonify({
                    'error': 'Invalid file type. Only CSV and XLSX files are allowed.'
                })

            upload_data, error = s3_manager.generate_presigned_upload_url(
                filename,
                content_type
            )
            if error:
                return jsonify({'error': f'Could not generate upload URL: {error}'})

            return jsonify({
                'success': True,
                'upload_url': upload_data['url'],
                'upload_fields': upload_data['fields'],
                's3_key': upload_data['key']
            })
        except Exception as e:
            logger.error(f"❌ Upload URL error: {e}")
            return jsonify({'error': f'Upload URL service error: {str(e)}'})

    @app.route('/admin/process-upload', methods=['POST'])
    @login_required
    def process_upload():
        """Upload processing with RDS integration (optimized + streaming for CSV)."""
        start_time = time.time()

        def log_step(message: str):
            elapsed = time.time() - start_time
            logger.info(f"[process_upload] [{elapsed:6.2f}s] {message}")

        def detect_csv_encoding(file_bytes):
            """Try a few encodings quickly on a small sample."""
            for encoding in ['utf-8', 'utf-8-sig', 'latin-1']:
                try:
                    _ = pd.read_csv(
                        BytesIO(file_bytes),
                        dtype=PANDAS_OPTIONS['dtype'],
                        low_memory=False,
                        usecols=REQUIRED_COLUMNS,
                        encoding=encoding,
                        nrows=100
                    )
                    return encoding
                except UnicodeDecodeError:
                    continue
                except Exception:
                    # ignore non-encoding errors here
                    continue
            return None

        try:
            log_step("Request received")
            data = request.get_json()
            if not data:
                log_step("Invalid request format (no JSON body)")
                return jsonify({'error': 'Invalid request format'})

            s3_key = data.get('s3_key', '').strip()
            mode = data.get('upload_type', 'replace').strip()

            if not s3_key:
                log_step("Missing S3 key in request")
                return jsonify({'error': 'S3 key is required'})
            if mode not in ['replace', 'append']:
                log_step(f"Invalid upload mode: {mode}")
                return jsonify({'error': 'Invalid upload mode'})

            logger.info(f"📥 Processing upload: {s3_key} (mode: {mode})")
            log_step(f"Validated request: s3_key={s3_key}, mode={mode}")

            # --- 1. Download file from S3 ---
            file_content, error = s3_manager.read_file_from_s3(s3_key)
            if error:
                log_step(f"Error reading file from S3: {error}")
                return jsonify({'error': f'Could not read uploaded file: {error}'})
            log_step("Downloaded file from S3")

            file_extension = s3_key.lower().split('.')[-1]

            # Preload existing dedup keys once for append mode
            existing_keys = get_existing_dedup_keys() if mode == 'append' else set()
            if mode == 'append':
                log_step(f"Loaded {len(existing_keys)} existing dedup keys for append mode")

            # --- 2. Backup existing S3 snapshot BEFORE processing ---
            backup_key = backup_current_data_in_s3(mode)
            log_step(f"Backup completed (key={backup_key})")

            # --- XLSX path: non-streaming but optimized --- #
            if file_extension == 'xlsx':
                try:
                    df = pd.read_excel(
                        BytesIO(file_content),
                        engine='openpyxl',
                        usecols=REQUIRED_COLUMNS
                    )
                    original_rows = len(df)
                    log_step(f"Parsed XLSX into DataFrame with {original_rows} rows")
                except Exception as read_error:
                    logger.error(f"❌ Error reading XLSX file: {read_error}")
                    s3_manager.delete_file(s3_key)
                    log_step(f"Error while parsing uploaded XLSX file: {read_error}")
                    return jsonify({'error': 'Invalid file format or corrupted file'})

                if df is None or df.empty:
                    s3_manager.delete_file(s3_key)
                    log_step("Uploaded XLSX contained no data (empty DataFrame)")
                    return jsonify({'error': 'Uploaded file contains no data'})

                try:
                    df_cleaned = clean_employee_dataframe_chunked(df)
                    if df_cleaned is None or df_cleaned.empty:
                        s3_manager.delete_file(s3_key)
                        log_step("Cleaning XLSX left no valid rows")
                        return jsonify({
                            'error': 'No valid data remains after cleaning'
                        })
                    cleaned_rows = len(df_cleaned)
                    dropped_rows = original_rows - cleaned_rows
                    log_step(
                        f"Cleaned XLSX DataFrame: {cleaned_rows} rows "
                        f"(dropped {dropped_rows} invalid rows)"
                    )
                except Exception as cleaning_error:
                    logger.error(f"❌ Error cleaning XLSX data: {cleaning_error}")
                    s3_manager.delete_file(s3_key)
                    log_step(f"Data cleaning failed: {cleaning_error}")
                    return jsonify({
                        'error': f'Data cleaning failed: {str(cleaning_error)}'
                    })

                missing_columns = [
                    col for col in REQUIRED_COLUMNS if col not in df_cleaned.columns
                ]
                if missing_columns:
                    s3_manager.delete_file(s3_key)
                    log_step(f"Missing required columns: {', '.join(missing_columns)}")
                    return jsonify({
                        'error': f'Missing required columns: {", ".join(missing_columns)}'
                    })

                dedup_keys_in_file: set[str] = set()
                df_cleaned, duplicates_in_file, dedup_keys_in_file = _deduplicate_cleaned_dataframe(
                    df_cleaned,
                    dedup_keys_in_file
                )
                if df_cleaned is None or df_cleaned.empty:
                    s3_manager.delete_file(s3_key)
                    log_step("Uploaded XLSX contained only duplicate rows after deduplication")
                    return jsonify({
                        'error': 'Uploaded file contains only duplicate rows'
                    })

                cleaned_rows = int(len(df_cleaned))
                dropped_rows = original_rows - cleaned_rows
                duplicates_skipped = 0
                duplicates_total = int(duplicates_in_file)
                s3_error = None
                certificate_stats = {'created': 0, 'skipped': 0, 'failed': 0}
                certificate_error = None

                if mode == 'append':
                    df_for_db, duplicates_skipped = filter_new_records_for_append(df_cleaned, existing_keys)
                    duplicates_skipped = int(duplicates_skipped)
                    duplicates_total = int(duplicates_in_file + duplicates_skipped)
                    cleaned_rows = int(len(df_for_db))
                    dropped_rows = original_rows - cleaned_rows
                    log_step(
                        f"Append dedup complete: {duplicates_total} duplicates skipped, "
                        f"{cleaned_rows} new rows to insert"
                    )

                    if cleaned_rows == 0:
                        s3_manager.delete_file(s3_key)
                        invalidate_stats_cache()
                        log_step("No new rows to append after deduplication")
                        return jsonify({
                            'success': True,
                            'message': 'No new records to append; all rows were duplicates.',
                            'stats': {
                                'mode': 'append',
                                'count': 0,
                                'added': 0,
                                'dropped': int(dropped_rows),
                                'invalid': int(dropped_rows),
                                'duplicates': duplicates_total,
                                'skipped': duplicates_total,
                                'total': int(dropped_rows + duplicates_total),
                                'rds_success': True,
                                's3_success': True
                            }
                        })

                    rds_result, rds_error = save_cleaned_data_chunk_to_rds(
                        df_for_db,
                        already_mapped=True
                    )
                    log_step(
                        f"save_cleaned_data_chunk_to_rds completed "
                        f"(success={rds_result}, error={rds_error})"
                    )

                    if rds_result:
                        certificate_stats, certificate_error = generate_certificates_for_dataframe(
                            df_for_db,
                            already_mapped=True
                        )
                        log_step(
                            "Certificate generation completed for append XLSX "
                            f"(created={certificate_stats['created']}, "
                            f"skipped={certificate_stats['skipped']}, "
                            f"failed={certificate_stats['failed']}, "
                            f"error={certificate_error})"
                        )
                        s3_result, s3_error = export_rds_to_csv_snapshot()
                        log_step(
                            f"export_rds_to_csv_snapshot completed "
                            f"(success={s3_result}, error={s3_error})"
                        )
                    else:
                        s3_result = False
                else:
                    s3_result = write_employee_data_to_s3(df_cleaned)
                    log_step(f"write_employee_data_to_s3 completed (success={s3_result})")

                    rds_result, rds_error = save_cleaned_data_to_rds(df_cleaned)
                    log_step(
                        f"save_cleaned_data_to_rds completed "
                        f"(success={rds_result}, error={rds_error})"
                    )
                    if rds_result:
                        certificate_stats, certificate_error = generate_certificates_for_dataframe(df_cleaned)
                        log_step(
                            "Certificate generation completed for replace XLSX "
                            f"(created={certificate_stats['created']}, "
                            f"skipped={certificate_stats['skipped']}, "
                            f"failed={certificate_stats['failed']}, "
                            f"error={certificate_error})"
                        )

                # Build enhanced KB from RDS
                kb_content, record_count, issuer_count, kb_error = create_enhanced_knowledge_base_from_rds()
                if kb_content and not kb_error:
                    kb_uploaded = s3_manager.upload_text_to_s3(
                        kb_content,
                        'knowledge-base/employee_data.txt'
                    )
                    if kb_uploaded:
                        logger.info("Knowledge base (enhanced) created successfully from RDS (XLSX path)")
                        log_step("KB created & uploaded to S3 (XLSX path)")
                    else:
                        log_step("KB generation OK, but S3 upload failed (XLSX path)")
                else:
                    log_step(f"KB generation from RDS failed (XLSX path): {kb_error}")

                if s3_result and rds_result:
                    if mode == 'append':
                        success_msg = (
                            f'Data appended successfully! {cleaned_rows} new records added '
                            f'(duplicates skipped: {duplicates_total})'
                        )
                    else:
                        success_msg = (
                            f'Data replaced successfully! {cleaned_rows} records loaded '
                            f'(S3 backup and RDS updated)'
                        )
                elif rds_result:
                    success_msg = (
                        'Data saved in RDS successfully, but S3 backup failed'
                    )
                elif s3_result:
                    success_msg = (
                        f'Data saved in S3, but database update failed: {rds_error or s3_error}'
                    )
                else:
                    s3_manager.delete_file(s3_key)
                    raise Exception("Failed to write data to both S3 and RDS")

                if dropped_rows > 0:
                    success_msg += f', {dropped_rows} invalid records cleaned'
                if mode == 'append' and duplicates_total > 0:
                    success_msg += f', {duplicates_total} duplicates skipped'
                elif mode == 'replace' and duplicates_total > 0:
                    success_msg += f', {duplicates_total} duplicate rows removed'

                result_stats = {
                    'mode': mode,
                    'count': int(cleaned_rows),
                    'added': int(cleaned_rows),
                    'dropped': int(dropped_rows),
                    'invalid': int(dropped_rows),
                    'duplicates': int(duplicates_total),
                    'skipped': int(duplicates_total),
                    'total': int(cleaned_rows + duplicates_total + max(dropped_rows, 0)),
                    'rds_success': bool(rds_result),
                    's3_success': bool(s3_result)
                }
                if certificate_stats:
                    result_stats['certificates'] = {
                        'created': int(certificate_stats.get('created', 0)),
                        'skipped': int(certificate_stats.get('skipped', 0)),
                        'failed': int(certificate_stats.get('failed', 0)),
                        'error': certificate_error
                    }

                s3_manager.delete_file(s3_key)
                invalidate_stats_cache()
                del df, df_cleaned, file_content
                gc.collect()
                log_step("Cleanup complete and stats cache invalidated (XLSX)")

                logger.info(f"✅ Upload (XLSX) processed successfully: {success_msg}")
                log_step("Finished process_upload for XLSX successfully (returning 200)")
                return jsonify({
                    'success': True,
                    'message': success_msg,
                    'stats': result_stats
                })

            # --- CSV path: streaming in chunks --- #
            encoding = detect_csv_encoding(file_content)
            if not encoding:
                s3_manager.delete_file(s3_key)
                log_step("Failed to detect CSV encoding")
                return jsonify({'error': 'Could not decode CSV file'})

            log_step(f"Detected CSV encoding: {encoding}")

            if mode == 'replace':
                # Clear RDS table once for replace mode
                clear_ok, clear_error = clear_employee_data_in_rds()
                if not clear_ok:
                    s3_manager.delete_file(s3_key)
                    log_step(f"Failed to clear RDS table: {clear_error}")
                    return jsonify({
                        'error': f'Upload processing failed while clearing database: {clear_error}'
                    })

                total_original_rows = 0
                total_cleaned_rows = 0
                duplicates_in_file = 0
                required_checked = False
                dedup_keys_in_file: set[str] = set()
                certificate_stats = {'created': 0, 'skipped': 0, 'failed': 0}
                certificate_error = None

                # Temp file for building CSV for S3 snapshot (streaming)
                temp_file = tempfile.NamedTemporaryFile(
                    mode='w+', newline='', suffix='.csv', delete=False, encoding='utf-8'
                )
                temp_file_path = temp_file.name
                wrote_header = False

                try:
                    csv_stream = BytesIO(file_content)
                    for chunk in pd.read_csv(
                        csv_stream,
                        dtype=PANDAS_OPTIONS['dtype'],
                        low_memory=False,
                        usecols=REQUIRED_COLUMNS,
                        encoding=encoding,
                        chunksize=CHUNK_SIZE
                    ):
                        chunk_original_rows = len(chunk)
                        total_original_rows += chunk_original_rows

                        cleaned_chunk = clean_employee_dataframe_chunked(chunk)
                        if cleaned_chunk is None or cleaned_chunk.empty:
                            continue

                        if not required_checked:
                            missing_columns = [
                                col for col in REQUIRED_COLUMNS if col not in cleaned_chunk.columns
                            ]
                            if missing_columns:
                                s3_manager.delete_file(s3_key)
                                log_step(f"Missing required columns in first chunk: {', '.join(missing_columns)}")
                                temp_file.close()
                                os.unlink(temp_file_path)
                                return jsonify({
                                    'error': f'Missing required columns: {", ".join(missing_columns)}'
                                })
                            required_checked = True

                        cleaned_chunk, chunk_duplicates, dedup_keys_in_file = _deduplicate_cleaned_dataframe(
                            cleaned_chunk,
                            dedup_keys_in_file
                        )
                        duplicates_in_file += chunk_duplicates
                        if cleaned_chunk is None or cleaned_chunk.empty:
                            continue

                        chunk_cleaned_rows = len(cleaned_chunk)
                        total_cleaned_rows += chunk_cleaned_rows

                        # Append to RDS
                        rds_result, rds_error = save_cleaned_data_chunk_to_rds(cleaned_chunk)
                        if not rds_result:
                            s3_manager.delete_file(s3_key)
                            log_step(f"RDS chunk save failed: {rds_error}")
                            temp_file.close()
                            os.unlink(temp_file_path)
                            return jsonify({
                                'error': f'Upload processing failed while saving to database: {rds_error}'
                            })
                        chunk_cert_stats, chunk_cert_error = generate_certificates_for_dataframe(cleaned_chunk)
                        certificate_stats['created'] += chunk_cert_stats.get('created', 0)
                        certificate_stats['skipped'] += chunk_cert_stats.get('skipped', 0)
                        certificate_stats['failed'] += chunk_cert_stats.get('failed', 0)
                        if chunk_cert_error and not certificate_error:
                            certificate_error = chunk_cert_error

                        # Append to temp CSV for S3 snapshot
                        cleaned_chunk.to_csv(
                            temp_file,
                            index=False,
                            header=not wrote_header,
                            encoding='utf-8'
                        )
                        wrote_header = True

                    temp_file.flush()

                    if total_cleaned_rows == 0:
                        s3_manager.delete_file(s3_key)
                        log_step("No valid data in any CSV chunk after cleaning")
                        temp_file.close()
                        os.unlink(temp_file_path)
                        return jsonify({
                            'error': 'No valid data remains after cleaning'
                        })

                    dropped_rows = total_original_rows - total_cleaned_rows
                    log_step(
                        f"Finished CSV chunk processing: {total_cleaned_rows} cleaned rows "
                        f"(dropped {dropped_rows})"
                    )

                    # --- Upload combined CSV from temp file to S3 snapshot --- #
                    temp_file.seek(0)
                    csv_content = temp_file.read()
                    temp_file.close()
                    s3_result = write_employee_csv_string_to_s3(csv_content)
                    log_step(f"write_employee_csv_string_to_s3 completed (success={s3_result})")

                    # Build enhanced KB from RDS
                    kb_content, record_count, issuer_count, kb_error = create_enhanced_knowledge_base_from_rds()
                    if kb_content and not kb_error:
                        kb_uploaded = s3_manager.upload_text_to_s3(
                            kb_content,
                            'knowledge-base/employee_data.txt'
                        )
                        if kb_uploaded:
                            logger.info("Knowledge base (enhanced) created successfully from RDS (CSV path)")
                            log_step("KB created & uploaded to S3 (CSV path)")
                        else:
                            log_step("KB generation OK, but S3 upload failed (CSV path)")
                    else:
                        log_step(f"KB generation from RDS failed (CSV path): {kb_error}")

                    if not s3_result:
                        success_msg = (
                            'Data replaced in RDS successfully, but S3 backup failed'
                        )
                    else:
                        success_msg = (
                            f'Data replaced successfully! {total_cleaned_rows} records loaded '
                            f'(S3 backup and RDS updated)'
                        )

                    if dropped_rows > 0:
                        success_msg += f', {dropped_rows} invalid records cleaned'
                    if duplicates_in_file > 0:
                        success_msg += f', {duplicates_in_file} duplicate rows removed'

                    result_stats = {
                        'mode': 'replace',
                        'count': int(total_cleaned_rows),
                        'added': int(total_cleaned_rows),
                        'dropped': int(dropped_rows),
                        'invalid': int(dropped_rows),
                        'duplicates': int(duplicates_in_file),
                        'skipped': int(duplicates_in_file),
                        'total': int(total_cleaned_rows + duplicates_in_file + max(dropped_rows, 0)),
                        'rds_success': True,
                        's3_success': bool(s3_result)
                    }
                    result_stats['certificates'] = {
                        'created': int(certificate_stats.get('created', 0)),
                        'skipped': int(certificate_stats.get('skipped', 0)),
                        'failed': int(certificate_stats.get('failed', 0)),
                        'error': certificate_error
                    }

                    # Cleanup
                    s3_manager.delete_file(s3_key)
                    invalidate_stats_cache()
                    del file_content
                    gc.collect()
                    os.unlink(temp_file_path)
                    log_step("Cleanup complete and stats cache invalidated (CSV)")

                    logger.info(f"✅ Upload (CSV) processed successfully: {success_msg}")
                    log_step("Finished process_upload for CSV successfully (returning 200)")
                    return jsonify({
                        'success': True,
                        'message': success_msg,
                        'stats': result_stats
                    })
                except Exception as processing_error:
                    try:
                        temp_file.close()
                        if os.path.exists(temp_file_path):
                            os.unlink(temp_file_path)
                    except Exception:
                        pass
                    s3_manager.delete_file(s3_key)
                    logger.error(f"❌ Processing error (CSV): {processing_error}")
                    log_step(f"Processing error occurred (CSV): {processing_error}")
                    return jsonify({
                        'error': f'Upload processing failed: {str(processing_error)}'
                    })
            else:
                total_original_rows = 0
                total_cleaned_rows = 0
                total_duplicates = 0
                required_checked = False
                certificate_stats = {'created': 0, 'skipped': 0, 'failed': 0}
                certificate_error = None

                try:
                    csv_stream = BytesIO(file_content)
                    for chunk in pd.read_csv(
                        csv_stream,
                        dtype=PANDAS_OPTIONS['dtype'],
                        low_memory=False,
                        usecols=REQUIRED_COLUMNS,
                        encoding=encoding,
                        chunksize=CHUNK_SIZE
                    ):
                        chunk_original_rows = len(chunk)
                        total_original_rows += chunk_original_rows

                        cleaned_chunk = clean_employee_dataframe_chunked(chunk)
                        if cleaned_chunk is None or cleaned_chunk.empty:
                            continue

                        if not required_checked:
                            missing_columns = [
                                col for col in REQUIRED_COLUMNS if col not in cleaned_chunk.columns
                            ]
                            if missing_columns:
                                s3_manager.delete_file(s3_key)
                                log_step(f"Missing required columns in first chunk: {', '.join(missing_columns)}")
                                return jsonify({
                                    'error': f'Missing required columns: {", ".join(missing_columns)}'
                                })
                            required_checked = True

                        deduped_chunk, dupes = filter_new_records_for_append(cleaned_chunk, existing_keys)
                        if deduped_chunk is None or deduped_chunk.empty:
                            total_duplicates += dupes
                            continue

                        chunk_cleaned_rows = len(deduped_chunk)
                        total_cleaned_rows += chunk_cleaned_rows
                        total_duplicates += dupes

                        rds_result, rds_error = save_cleaned_data_chunk_to_rds(
                            deduped_chunk,
                            already_mapped=True
                        )
                        if not rds_result:
                            s3_manager.delete_file(s3_key)
                            log_step(f"RDS chunk save failed (append): {rds_error}")
                            return jsonify({
                                'error': f'Upload processing failed while saving to database: {rds_error}'
                            })
                        chunk_cert_stats, chunk_cert_error = generate_certificates_for_dataframe(
                            deduped_chunk,
                            already_mapped=True
                        )
                        certificate_stats['created'] += chunk_cert_stats.get('created', 0)
                        certificate_stats['skipped'] += chunk_cert_stats.get('skipped', 0)
                        certificate_stats['failed'] += chunk_cert_stats.get('failed', 0)
                        if chunk_cert_error and not certificate_error:
                            certificate_error = chunk_cert_error

                    if total_cleaned_rows == 0:
                        s3_manager.delete_file(s3_key)
                        log_step("No new rows to append after deduplication")
                        return jsonify({
                            'success': True,
                            'message': 'No new records to append; all rows were duplicates.',
                            'stats': {
                                'mode': 'append',
                                'count': 0,
                                'added': 0,
                                'dropped': int(total_original_rows),
                                'invalid': int(total_original_rows),
                                'duplicates': int(total_duplicates),
                                'skipped': int(total_duplicates),
                                'total': int(total_original_rows),
                                'rds_success': True,
                                's3_success': True
                            }
                        })

                    dropped_rows = int(total_original_rows - total_cleaned_rows)
                    log_step(
                        f"Finished CSV append processing: {total_cleaned_rows} new rows, "
                        f"{total_duplicates} duplicates skipped"
                    )

                    s3_result, s3_error = export_rds_to_csv_snapshot()
                    log_step(
                        f"export_rds_to_csv_snapshot completed "
                        f"(success={s3_result}, error={s3_error})"
                    )

                    kb_content, record_count, issuer_count, kb_error = create_enhanced_knowledge_base_from_rds()
                    if kb_content and not kb_error:
                        kb_uploaded = s3_manager.upload_text_to_s3(
                            kb_content,
                            'knowledge-base/employee_data.txt'
                        )
                        if kb_uploaded:
                            logger.info("Knowledge base (enhanced) created successfully from RDS (CSV append path)")
                            log_step("KB created & uploaded to S3 (CSV append path)")
                        else:
                            log_step("KB generation OK, but S3 upload failed (CSV append path)")
                    else:
                        log_step(f"KB generation from RDS failed (CSV append path): {kb_error}")

                    success_msg = (
                        f'Data appended successfully! {total_cleaned_rows} new records added '
                        f'(duplicates skipped: {total_duplicates})'
                    )
                    if dropped_rows > 0:
                        success_msg += f', {dropped_rows} invalid records cleaned'

                    result_stats = {
                        'mode': 'append',
                        'count': int(total_cleaned_rows),
                        'added': int(total_cleaned_rows),
                        'dropped': int(dropped_rows),
                        'invalid': int(dropped_rows),
                        'duplicates': int(total_duplicates),
                        'skipped': int(total_duplicates),
                        'total': int(total_cleaned_rows + total_duplicates + max(dropped_rows, 0)),
                        'rds_success': True,
                        's3_success': bool(s3_result)
                    }
                    result_stats['certificates'] = {
                        'created': int(certificate_stats.get('created', 0)),
                        'skipped': int(certificate_stats.get('skipped', 0)),
                        'failed': int(certificate_stats.get('failed', 0)),
                        'error': certificate_error
                    }

                    s3_manager.delete_file(s3_key)
                    invalidate_stats_cache()
                    del file_content
                    gc.collect()
                    log_step("Cleanup complete and stats cache invalidated (CSV append)")

                    logger.info(f"✅ Upload (CSV append) processed successfully: {success_msg}")
                    log_step("Finished process_upload for CSV append successfully (returning 200)")
                    return jsonify({
                        'success': True,
                        'message': success_msg,
                        'stats': result_stats
                    })
                except Exception as processing_error:
                    s3_manager.delete_file(s3_key)
                    logger.error(f"❌ Processing error (CSV append): {processing_error}")
                    log_step(f"Processing error occurred (CSV append): {processing_error}")
                    return jsonify({
                        'error': f'Upload processing failed: {str(processing_error)}'
                    })
        except Exception as e:
            logger.error(f"❌ Upload processing error: {e}")
            log_step(f"Top-level upload processing error: {e}")
            return jsonify({
                'error': f'Upload processing failed: {str(e)}'
            })

    # --- Refresh Knowledge Base from DB (using shared service) --- #

    @app.route('/admin/refresh-knowledge-base', methods=['POST'])
    @login_required
    def refresh_knowledge_base():
        """Create enhanced knowledge base with issuer-level-qualifier breakdown."""
        try:
            logger.info("📚 Creating enhanced knowledge base (manual refresh)...")
            kb_content, record_count, issuer_count, error = create_enhanced_knowledge_base_from_rds()
            if error or not kb_content:
                return jsonify({
                    "success": False,
                    "error": error or "Failed to create knowledge base"
                })

            kb_uploaded = s3_manager.upload_text_to_s3(
                kb_content,
                'knowledge-base/employee_data.txt'
            )
            if kb_uploaded:
                logger.info("✅ Knowledge base uploaded successfully (manual refresh)")
                return jsonify({
                    "success": True,
                    "message": (
                        f"Enhanced knowledge base created with {record_count} records "
                        f"covering {issuer_count} issuers"
                    )
                })
            return jsonify({"success": False, "error": "S3 upload failed"})
        except Exception as e:
            logger.error(f"KB error: {e}")
            return jsonify({"success": False, "error": str(e)})

    @app.route('/admin/generate-certificates', methods=['POST'])
    @login_required
    def generate_certificates():
        """Generate certificates for every record in the database."""
        if not certificate_generation_lock.acquire(blocking=False):
            return jsonify({
                "success": False,
                "error": "Certificate generation is already in progress. Please try again later.",
            }), 409

        certificate_generation_state.update({
            "running": True,
            "last_error": None,
            "last_finished_at": None,
        })

        def run_certificate_generation():
            try:
                with app.app_context():
                    logger.info("🎓 Generating certificates for all records in RDS (async)...")
                    counts, error = generate_certificates_for_rds()
                    certificate_generation_state["last_counts"] = counts
                    certificate_generation_state["last_error"] = error
                    certificate_generation_state["last_finished_at"] = datetime.utcnow().isoformat()
                    if error:
                        logger.warning(f"Certificate generation completed with error: {error}")
                    else:
                        logger.info(
                            "✅ Certificate generation completed: %s created, %s skipped, %s failed.",
                            counts["created"],
                            counts["skipped"],
                            counts["failed"],
                        )
            except Exception as e:
                logger.error(f"Certificate generation error: {e}")
                certificate_generation_state["last_error"] = str(e)
            finally:
                certificate_generation_state["running"] = False
                certificate_generation_lock.release()

        threading.Thread(target=run_certificate_generation, daemon=True).start()

        return jsonify({
            "success": True,
            "message": "Certificate generation started. Check back later for results.",
            "status": "started",
        }), 202

    # --- Stats & chart APIs --- #

    @app.route('/admin/stats')
    @login_required
    def get_stats():
        """Get cached statistics from RDS."""
        try:
            stats = get_employee_stats_cached()
            return jsonify(stats)
        except Exception as e:
            logger.error(f"❌ Stats error: {e}")
            return jsonify({
                'error': f'Error reading stats: {str(e)}',
                'total_employees': 0,
                'completed_journey': 0,
                'in_progress': 0,
                'all_records': 0,
                'rds_records': 0
            })

    @app.route('/admin/chart-data')
    @login_required
    def chart_data():
        """RDS-optimized chart data."""
        try:
            issuer = request.args.get('issuer', '').strip()
            chart_data_result = get_chart_data_from_rds(issuer)
            return jsonify(chart_data_result)
        except Exception as e:
            logger.error(f"❌ Chart data error: {e}")
            return jsonify({
                'error': f'Error reading chart data: {str(e)}',
                'labels': ['Level 0', 'Level 1', 'Level 2'],
                'values': [0, 0, 0],
                'issuers': []
            })

    @app.route('/admin/certification-data')
    @login_required
    def get_certification_data():
        """Get certification data by level and issuer."""
        try:
            issuer = request.args.get('issuer', '').strip()
            query = db.session.query(
                EmployeeRecord.level,
                EmployeeRecord.qualifier,
                func.count(EmployeeRecord.id)
            )
            if issuer:
                query = query.filter(EmployeeRecord.issuer == issuer)
            results = query.group_by(
                EmployeeRecord.level,
                EmployeeRecord.qualifier
            ).all()
            levels = ['Level 0', 'Level 1', 'Level 2', 'Level 3']
            certified = {level: 0 for level in levels}
            trained = {level: 0 for level in levels}
            for level, qualifier, count in results:
                if level in certified:
                    q = (qualifier or "").lower()
                    if 'certified' in q:
                        certified[level] += count
                    elif 'trained' in q:
                        trained[level] += count
            all_issuers = db.session.query(
                func.distinct(EmployeeRecord.issuer)
            ).filter(
                EmployeeRecord.issuer.isnot(None)
            ).all()
            issuers = sorted([i[0] for i in all_issuers if i[0]])
            return jsonify({
                "labels": levels,
                "certified": [certified[l] for l in levels],
                "trained": [trained[l] for l in levels],
                "issuers": issuers,
                "selected_issuer": issuer
            })
        except Exception as e:
            logger.error(f"Certification data error: {e}")
            return jsonify({"error": str(e)})

    @app.route('/admin/exports', methods=['GET'])
    @login_required
    def export_filtered_records():
        """Export filtered employee records as XLSX or PPTX."""
        try:
            export_format = request.args.get('format', 'xlsx').lower()
            filters = _normalize_export_filters(request.args)
            query = _apply_export_filters(db.session.query(EmployeeRecord), filters)
            records = query.all()
            records = _filter_records_by_date(records, filters)

            if not records:
                return jsonify({"success": False, "error": "No records found for the selected filters"}), 404

            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            if export_format == 'pptx':
                pptx_stream = _build_pptx_report(records, filters)
                filename = f"ai_academy_insights_{timestamp}.pptx"
                return send_file(
                    pptx_stream,
                    mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    as_attachment=True,
                    download_name=filename
                )

            export_rows = _records_to_export_rows(records)
            df = pd.DataFrame(export_rows)
            output = BytesIO()
            df.to_excel(output, index=False)
            output.seek(0)
            filename = f"ai_academy_export_{timestamp}.xlsx"
            return send_file(
                output,
                mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                as_attachment=True,
                download_name=filename
            )
        except Exception as e:
            logger.error(f"❌ Export error: {e}")
            return jsonify({"success": False, "error": str(e)}), 500
